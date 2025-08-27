"""
PPTX Rendering MCP Tool
Generates executive roadmap presentations using python-pptx with citations.
"""
import sys
sys.path.append("/app")
import os
import tempfile
import logging
from typing import Dict, Any, List, Optional
from datetime import datetime, timezone
import json
import base64

# PPTX generation imports
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.dml.color import RGBColor
    PPTX_DEPS_AVAILABLE = True
except ImportError:
    PPTX_DEPS_AVAILABLE = False

from api.security import SecurityPolicy
from config import MCPConfig

logger = logging.getLogger(__name__)

class PPTXRenderTool:
    """
    MCP tool for generating executive PPTX presentations.
    
    Features:
    - Executive roadmap deck generation
    - Citations and source tracking
    - Professional templates
    - Configurable branding
    - Export to base64 or file path
    """
    
    TOOL_NAME = "pptx.render"
    
    # Default slide templates
    SLIDE_TEMPLATES = {
        "title": {
            "layout_index": 0,
            "placeholders": {
                "title": 0,
                "subtitle": 1
            }
        },
        "content": {
            "layout_index": 1,
            "placeholders": {
                "title": 0,
                "content": 1
            }
        },
        "two_content": {
            "layout_index": 3,
            "placeholders": {
                "title": 0,
                "left_content": 1,
                "right_content": 2
            }
        },
        "blank": {
            "layout_index": 6,
            "placeholders": {}
        }
    }
    
    # Default color scheme
    DEFAULT_COLORS = {
        "primary": (0, 102, 204),      # Blue
        "secondary": (102, 102, 102),   # Gray
        "accent": (255, 102, 0),        # Orange
        "text": (51, 51, 51),           # Dark gray
        "background": (255, 255, 255)   # White
    }
    
    def __init__(self, config: MCPConfig):
        """Initialize PPTX renderer with configuration."""
        self.config = config
        self.security = SecurityPolicy(config)
        
        if not PPTX_DEPS_AVAILABLE:
            logger.warning("PPTX dependencies not available. Tool will run in mock mode.")
    
    def validate_presentation_data(self, payload: Dict[str, Any]) -> Dict[str, Any]:
        """
        Validate presentation data structure.
        
        Args:
            payload: Tool payload containing presentation data
            
        Returns:
            Dict containing validated presentation configuration
        """
        if "presentation" not in payload:
            raise ValueError("Missing required field: presentation")
        
        presentation_data = payload["presentation"]
        
        # Required fields
        required_fields = ["title", "slides"]
        for field in required_fields:
            if field not in presentation_data:
                raise ValueError(f"Missing required presentation field: {field}")
        
        # Validate slides structure
        slides = presentation_data["slides"]
        if not isinstance(slides, list) or len(slides) == 0:
            raise ValueError("Presentation must contain at least one slide")
        
        for i, slide in enumerate(slides):
            if not isinstance(slide, dict):
                raise ValueError(f"Slide {i} must be a dictionary")
            
            if "type" not in slide:
                raise ValueError(f"Slide {i} missing required field: type")
            
            slide_type = slide["type"]
            if slide_type not in self.SLIDE_TEMPLATES:
                raise ValueError(f"Slide {i} has unsupported type: {slide_type}")
        
        # Validate optional configuration
        config = {
            "title": presentation_data["title"],
            "subtitle": presentation_data.get("subtitle", ""),
            "author": presentation_data.get("author", "AI Cyber Maturity Assessment"),
            "slides": slides,
            "branding": presentation_data.get("branding", {}),
            "citations": presentation_data.get("citations", []),
            "template": presentation_data.get("template", "default")
        }
        
        return config
    
    def create_presentation(self, config: Dict[str, Any]) -> Any:
        """
        Create PPTX presentation object.
        
        Args:
            config: Validated presentation configuration
            
        Returns:
            PPTX Presentation object
        """
        if not PPTX_DEPS_AVAILABLE:
            return self._mock_presentation(config)
        
        # Create new presentation
        prs = Presentation()
        
        # Set presentation properties
        prs.core_properties.title = config["title"]
        prs.core_properties.author = config["author"]
        prs.core_properties.subject = "AI Cyber Maturity Assessment Roadmap"
        prs.core_properties.created = datetime.now(timezone.utc)
        
        return prs
    
    def add_title_slide(self, prs: Any, config: Dict[str, Any]) -> None:
        """Add title slide to presentation."""
        if not PPTX_DEPS_AVAILABLE:
            return
        
        # Use title slide layout
        title_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_layout)
        
        # Set title and subtitle
        title_placeholder = slide.shapes.title
        title_placeholder.text = config["title"]
        
        if slide.shapes.placeholders[1]:
            subtitle_placeholder = slide.shapes.placeholders[1]
            subtitle_text = config["subtitle"]
            if not subtitle_text:
                subtitle_text = f"Generated on {datetime.now().strftime('%B %d, %Y')}"
            subtitle_placeholder.text = subtitle_text
        
        # Apply branding colors if specified
        self._apply_slide_branding(slide, config.get("branding", {}))
    
    def add_content_slide(self, prs: Any, slide_data: Dict[str, Any], branding: Dict[str, Any]) -> None:
        """Add content slide to presentation."""
        if not PPTX_DEPS_AVAILABLE:
            return
        
        slide_type = slide_data["type"]
        template = self.SLIDE_TEMPLATES[slide_type]
        
        # Create slide with appropriate layout
        layout = prs.slide_layouts[template["layout_index"]]
        slide = prs.slides.add_slide(layout)
        
        # Set title
        if "title" in slide_data and slide.shapes.title:
            slide.shapes.title.text = slide_data["title"]
        
        # Add content based on slide type
        if slide_type == "content":
            self._add_bullet_content(slide, slide_data.get("content", []))
        elif slide_type == "two_content":
            self._add_two_column_content(slide, slide_data)
        
        # Apply branding
        self._apply_slide_branding(slide, branding)
    
    def _add_bullet_content(self, slide: Any, content_items: List[str]) -> None:
        """Add bullet point content to slide."""
        if not PPTX_DEPS_AVAILABLE or not content_items:
            return
        
        # Get content placeholder
        content_placeholder = None
        for placeholder in slide.shapes.placeholders:
            if placeholder.placeholder_format.idx == 1:  # Content placeholder
                content_placeholder = placeholder
                break
        
        if content_placeholder:
            text_frame = content_placeholder.text_frame
            text_frame.clear()
            
            for i, item in enumerate(content_items):
                if i == 0:
                    p = text_frame.paragraphs[0]
                else:
                    p = text_frame.add_paragraph()
                
                p.text = item
                p.level = 0
    
    def _add_two_column_content(self, slide: Any, slide_data: Dict[str, Any]) -> None:
        """Add two-column content to slide."""
        if not PPTX_DEPS_AVAILABLE:
            return
        
        left_content = slide_data.get("left_content", [])
        right_content = slide_data.get("right_content", [])
        
        # Add left column content
        if left_content and len(slide.shapes.placeholders) > 1:
            left_placeholder = slide.shapes.placeholders[1]
            if left_placeholder:
                text_frame = left_placeholder.text_frame
                text_frame.clear()
                for i, item in enumerate(left_content):
                    if i == 0:
                        p = text_frame.paragraphs[0]
                    else:
                        p = text_frame.add_paragraph()
                    p.text = item
        
        # Add right column content
        if right_content and len(slide.shapes.placeholders) > 2:
            right_placeholder = slide.shapes.placeholders[2]
            if right_placeholder:
                text_frame = right_placeholder.text_frame
                text_frame.clear()
                for i, item in enumerate(right_content):
                    if i == 0:
                        p = text_frame.paragraphs[0]
                    else:
                        p = text_frame.add_paragraph()
                    p.text = item
    
    def add_citations_slide(self, prs: Any, citations: List[Dict[str, Any]], branding: Dict[str, Any]) -> None:
        """Add citations slide to presentation."""
        if not PPTX_DEPS_AVAILABLE or not citations:
            return
        
        # Use content slide layout
        layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(layout)
        
        # Set title
        slide.shapes.title.text = "Sources and Citations"
        
        # Add citations content
        content_placeholder = slide.shapes.placeholders[1]
        text_frame = content_placeholder.text_frame
        text_frame.clear()
        
        for i, citation in enumerate(citations):
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            
            # Format citation
            title = citation.get("title", "Untitled")
            source = citation.get("source", "Unknown")
            date = citation.get("date", "")
            
            citation_text = f"{title} - {source}"
            if date:
                citation_text += f" ({date})"
            
            p.text = citation_text
            p.level = 0
        
        # Apply branding
        self._apply_slide_branding(slide, branding)
    
    def _apply_slide_branding(self, slide: Any, branding: Dict[str, Any]) -> None:
        """Apply branding colors and styles to slide."""
        if not PPTX_DEPS_AVAILABLE:
            return
        
        # Get colors from branding or use defaults
        colors = branding.get("colors", self.DEFAULT_COLORS)
        
        # Apply title formatting if title exists
        if slide.shapes.title:
            title_shape = slide.shapes.title
            if title_shape.text_frame:
                for paragraph in title_shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        font = run.font
                        font.size = Pt(32)
                        font.bold = True
                        if "primary" in colors:
                            font.color.rgb = RGBColor(*colors["primary"])
    
    def save_presentation(self, prs: Any, output_format: str = "base64") -> Dict[str, Any]:
        """
        Save presentation to specified format.
        
        Args:
            prs: PPTX Presentation object
            output_format: "base64" or "file"
            
        Returns:
            Dict containing saved presentation data
        """
        if not PPTX_DEPS_AVAILABLE:
            return self._mock_save_presentation(output_format)
        
        # Create temporary file
        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as temp_file:
            temp_path = temp_file.name
        
        try:
            # Save presentation
            prs.save(temp_path)
            
            if output_format == "base64":
                # Read file and encode to base64
                with open(temp_path, "rb") as f:
                    file_data = f.read()
                
                encoded_data = base64.b64encode(file_data).decode()
                
                return {
                    "format": "base64",
                    "data": encoded_data,
                    "filename": f"roadmap_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx",
                    "size_bytes": len(file_data)
                }
            
            elif output_format == "file":
                # Return file path
                file_size = os.path.getsize(temp_path)
                
                return {
                    "format": "file",
                    "path": temp_path,
                    "filename": os.path.basename(temp_path),
                    "size_bytes": file_size
                }
            
            else:
                raise ValueError(f"Unsupported output format: {output_format}")
        
        except Exception as e:
            # Clean up temp file on error
            try:
                os.unlink(temp_path)
            except:
                pass
            raise e
    
    def _mock_presentation(self, config: Dict[str, Any]) -> Dict[str, Any]:
        """Mock presentation for development/testing."""
        return {
            "mock_mode": True,
            "title": config["title"],
            "slide_count": len(config["slides"]) + 2,  # +2 for title and citations
            "config": config
        }
    
    def _mock_save_presentation(self, output_format: str) -> Dict[str, Any]:
        """Mock save presentation for development/testing."""
        mock_pptx_content = b"Mock PPTX file content for testing"
        
        if output_format == "base64":
            return {
                "format": "base64",
                "data": base64.b64encode(mock_pptx_content).decode(),
                "filename": f"mock_roadmap_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx",
                "size_bytes": len(mock_pptx_content),
                "mock_mode": True
            }
        else:
            # Create temporary mock file
            with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as temp_file:
                temp_file.write(mock_pptx_content)
                temp_path = temp_file.name
            
            return {
                "format": "file",
                "path": temp_path,
                "filename": os.path.basename(temp_path),
                "size_bytes": len(mock_pptx_content),
                "mock_mode": True
            }
    
    async def execute(self, payload: Dict[str, Any], engagement_id: str, call_id: str) -> Dict[str, Any]:
        """
        Execute PPTX rendering with full validation.
        
        Args:
            payload: Tool execution payload
            engagement_id: Engagement identifier for sandboxing
            call_id: Unique call identifier for tracking
            
        Returns:
            Dict containing rendered presentation and metadata
        """
        start_time = datetime.now(timezone.utc)
        
        try:
            # Validate presentation data
            config = self.validate_presentation_data(payload)
            
            # Get output preferences
            output_format = payload.get("output_format", "base64")
            if output_format not in ["base64", "file"]:
                raise ValueError("output_format must be 'base64' or 'file'")
            
            # Create presentation
            prs = self.create_presentation(config)
            
            # Add title slide
            self.add_title_slide(prs, config)
            
            # Add content slides
            for slide_data in config["slides"]:
                self.add_content_slide(prs, slide_data, config["branding"])
            
            # Add citations slide if citations provided
            if config["citations"]:
                self.add_citations_slide(prs, config["citations"], config["branding"])
            
            # Save presentation
            presentation_data = self.save_presentation(prs, output_format)
            
            # Calculate processing metrics
            processing_time = (datetime.now(timezone.utc) - start_time).total_seconds()
            
            # Prepare result
            result = {
                "success": True,
                "tool": self.TOOL_NAME,
                "call_id": call_id,
                "engagement_id": engagement_id,
                "presentation": presentation_data,
                "metadata": {
                    "title": config["title"],
                    "author": config["author"],
                    "slide_count": len(config["slides"]) + 1 + (1 if config["citations"] else 0),
                    "citations_count": len(config["citations"]),
                    "template": config["template"],
                    "has_branding": bool(config["branding"])
                },
                "processing_time_seconds": round(processing_time, 3),
                "timestamp": datetime.now(timezone.utc).isoformat()
            }
            
            # Log successful rendering
            logger.info(
                "PPTX presentation rendered successfully",
                extra={
                    "call_id": call_id,
                    "engagement_id": engagement_id,
                    "title": config["title"],
                    "slide_count": result["metadata"]["slide_count"],
                    "output_format": output_format,
                    "file_size_bytes": presentation_data.get("size_bytes", 0),
                    "processing_time_seconds": processing_time
                }
            )
            
            return result
            
        except Exception as e:
            error_result = {
                "success": False,
                "tool": self.TOOL_NAME,
                "call_id": call_id,
                "engagement_id": engagement_id,
                "error": str(e),
                "error_type": type(e).__name__,
                "timestamp": datetime.now(timezone.utc).isoformat()
            }
            
            logger.error(
                "PPTX rendering failed",
                extra={
                    "call_id": call_id,
                    "engagement_id": engagement_id,
                    "error": str(e),
                    "error_type": type(e).__name__
                }
            )
            
            return error_result


# Tool registration function
def register_tool(tool_registry: Dict[str, Any]) -> None:
    """Register the PPTX rendering tool with MCP gateway."""
    tool_registry[PPTXRenderTool.TOOL_NAME] = PPTXRenderTool